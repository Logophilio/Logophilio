import collections 
import collections.abc
import pptx

template = pptx.Presentation('template.pptx')
print(template.slides[0].shapes.element)
